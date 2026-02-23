"""
Conversione di tutti i tipi di file in formato PDF.

Pipeline a cascata (dal più qualitativo al fallback):
  DOCX/DOC  → 1) docx2pdf (Office COM)  2) mammoth+weasyprint  3) LibreOffice
  XLSX/XLS  → 1) docx2pdf (Office COM)  2) openpyxl+fpdf2       3) LibreOffice
  PPTX/PPT  → 1) docx2pdf (Office COM)  2) python-pptx+fpdf2    3) LibreOffice
  ODT/ODS   → 1) LibreOffice             2) mammoth (ODT)
  RTF/HTML  → 1) mammoth+weasyprint      2) LibreOffice
  TXT/CSV   → 1) fpdf2 (testo diretto)
  Immagini  → 1) img2pdf                 2) Pillow
  PDF       → copia diretta
  P7M       → estrazione + ricorsione

LibreOffice e' OPZIONALE: migliora la qualita' ma non e' indispensabile.
Se Microsoft Office e' installato, docx2pdf lo usa come prima scelta
(massima qualita', zero dipendenze extra).
"""

import os
import sys
import shutil
import subprocess
import tempfile
import uuid

# ── Estensioni per categoria ──────────────────────────────────────────────────

IMAGE_EXTENSIONS = {'.jpg', '.jpeg', '.png', '.gif', '.bmp', '.tiff', '.tif', '.webp'}

DOCX_EXTENSIONS  = {'.doc', '.docx', '.rtf'}
XLSX_EXTENSIONS  = {'.xls', '.xlsx', '.csv', '.ods'}
PPTX_EXTENSIONS  = {'.ppt', '.pptx', '.odp'}
ODT_EXTENSIONS   = {'.odt', '.odg'}
HTML_EXTENSIONS  = {'.html', '.htm'}
TEXT_EXTENSIONS  = {'.txt'}
XML_EXTENSIONS   = {'.xml'}

# ── Ricerca eseguibili di sistema ─────────────────────────────────────────────

def find_libreoffice() -> str:
    """Trova LibreOffice nel sistema. Ritorna il percorso o None."""
    import glob as _glob

    found = shutil.which('soffice')
    if found:
        return found

    if sys.platform == 'win32':
        patterns = [
            r'C:\Program Files\LibreOffice*\program\soffice.exe',
            r'C:\Program Files (x86)\LibreOffice*\program\soffice.exe',
        ]
        for p in patterns:
            matches = _glob.glob(p)
            if matches:
                return matches[-1]

    elif sys.platform == 'darwin':
        for p in [
            '/Applications/LibreOffice.app/Contents/MacOS/soffice',
            '/usr/local/bin/soffice',
            '/opt/homebrew/bin/soffice',
        ]:
            if os.path.isfile(p):
                return p

    return None


def has_microsoft_office() -> bool:
    """
    Verifica se Microsoft Office e' installato e utilizzabile da docx2pdf.
    """
    try:
        import docx2pdf  # noqa
    except ImportError:
        return False

    if sys.platform == 'win32':
        # Controlla se Word e' registrato come applicazione COM
        import winreg
        try:
            winreg.OpenKey(winreg.HKEY_LOCAL_MACHINE,
                           r'SOFTWARE\Microsoft\Office', 0, winreg.KEY_READ)
            return True
        except Exception:
            pass
        # Fallback: controlla percorsi comuni
        import glob as _glob
        patterns = [
            r'C:\Program Files\Microsoft Office\root\Office*\WINWORD.EXE',
            r'C:\Program Files (x86)\Microsoft Office\root\Office*\WINWORD.EXE',
            r'C:\Program Files\Microsoft Office\Office*\WINWORD.EXE',
        ]
        return any(_glob.glob(p) for p in patterns)

    elif sys.platform == 'darwin':
        return os.path.isdir('/Applications/Microsoft Word.app')

    return False


# ── Helper: LibreOffice headless ─────────────────────────────────────────────

def _convert_office_to_pdf(file_path: str, output_dir: str, lo_path: str) -> str:
    """Converti un file tramite LibreOffice headless."""
    tmp_dir = tempfile.mkdtemp(prefix='lo_conv_')
    try:
        tmp_input = os.path.join(tmp_dir, os.path.basename(file_path))
        shutil.copy2(file_path, tmp_input)

        subprocess.run(
            [lo_path, '--headless', '--norestore',
             '--convert-to', 'pdf', '--outdir', tmp_dir, tmp_input],
            capture_output=True, timeout=120,
            env={**os.environ, 'HOME': tmp_dir},
        )

        base = os.path.splitext(os.path.basename(file_path))[0]
        expected = os.path.join(tmp_dir, f'{base}.pdf')
        if not os.path.isfile(expected):
            pdfs = [f for f in os.listdir(tmp_dir) if f.endswith('.pdf')]
            if not pdfs:
                raise RuntimeError('LibreOffice non ha prodotto PDF')
            expected = os.path.join(tmp_dir, pdfs[0])

        dest = os.path.join(output_dir, f'{base}_{uuid.uuid4().hex[:8]}.pdf')
        shutil.move(expected, dest)
        return dest
    finally:
        shutil.rmtree(tmp_dir, ignore_errors=True)


# ── Helper: win32com diretto (Windows) ───────────────────────────────────────

def _try_win32com(docx_path: str, pdf_path: str) -> tuple:
    """
    Conversione tramite Word COM con soppressione completa dei dialoghi.
    Ritorna (percorso_pdf, None) oppure (None, str_errore).
    """
    try:
        import pythoncom
        from win32com import client as win32client
    except ImportError:
        return None, 'win32com non disponibile'

    word = None
    try:
        pythoncom.CoInitialize()
        word = win32client.Dispatch('Word.Application')
        word.Visible = False
        word.DisplayAlerts = 0   # wdAlertsNone: nessun dialogo

        doc = word.Documents.Open(
            os.path.abspath(docx_path),
            ReadOnly=True,
            AddToRecentFiles=False,
            ConfirmConversions=False,
        )
        try:
            doc.SaveAs2(os.path.abspath(pdf_path), FileFormat=17)  # wdFormatPDF
        finally:
            doc.Close(SaveChanges=0)   # wdDoNotSaveChanges

        if os.path.isfile(pdf_path) and os.path.getsize(pdf_path) > 0:
            return pdf_path, None
        return None, 'output PDF vuoto dopo SaveAs2'
    except Exception as e:
        return None, str(e)
    finally:
        if word is not None:
            try:
                word.Quit()
            except Exception:
                pass
        try:
            pythoncom.CoUninitialize()
        except Exception:
            pass


# ── Helper: AppleScript diretto per .doc su macOS ────────────────────────────

def _try_applescript_word(file_path: str, output_dir: str) -> tuple:
    """
    Apre il file in Microsoft Word via AppleScript e lo esporta come PDF.
    Funziona anche con .doc binario (Word 97-2003) che docx2pdf non gestisce.
    Ritorna (percorso_pdf, None) oppure (None, str_errore).
    """
    if sys.platform != 'darwin':
        return None, 'AppleScript disponibile solo su macOS'

    if not os.path.isdir('/Applications/Microsoft Word.app'):
        return None, 'Microsoft Word non trovato'

    base = os.path.splitext(os.path.basename(file_path))[0]
    dest = os.path.join(output_dir, f'{base}_{uuid.uuid4().hex[:8]}.pdf')
    abs_input = os.path.abspath(file_path)
    abs_output = os.path.abspath(dest)

    script = f'''
tell application "Microsoft Word"
    try
        open POSIX file "{abs_input}"
        delay 2
        set theDoc to active document
        save as theDoc file name "{abs_output}" file format format PDF
        close theDoc saving no
    on error errMsg
        return "ERRORE: " & errMsg
    end try
end tell
'''
    try:
        r = subprocess.run(
            ['osascript', '-e', script],
            capture_output=True, text=True, timeout=20
        )
        if os.path.isfile(dest) and os.path.getsize(dest) > 0:
            return dest, None
        err = (r.stderr or r.stdout or '').strip()
        return None, f'AppleScript Word: {err or "output vuoto"}'
    except Exception as e:
        return None, f'AppleScript Word: {e}'


# ── Helper: docx2pdf (Office COM / AppleScript) ───────────────────────────────

def _try_docx2pdf(file_path: str, output_dir: str) -> tuple:
    """
    Prova a convertire tramite Microsoft Office.
    Su Windows usa win32com direttamente (con dialog suppression).
    Su macOS usa docx2pdf (AppleScript).
    Ritorna (percorso_pdf, None) oppure (None, str_errore).
    """
    # Copia in una cartella trusted per evitare la Protected View di Word
    tmp_trusted = tempfile.mkdtemp(prefix='docx2pdf_')
    try:
        trusted_copy = os.path.join(tmp_trusted, os.path.basename(file_path))
        shutil.copy2(file_path, trusted_copy)

        # Rimuovi il flag Zone.Identifier (ADS "scaricato da internet")
        if sys.platform == 'win32':
            try:
                subprocess.run(
                    ['powershell', '-NonInteractive', '-WindowStyle', 'Hidden',
                     '-Command', f'Unblock-File -Path "{trusted_copy}"'],
                    capture_output=True, timeout=10
                )
            except Exception:
                pass

        base = os.path.splitext(os.path.basename(file_path))[0]
        dest = os.path.join(output_dir, f'{base}_{uuid.uuid4().hex[:8]}.pdf')

        if sys.platform == 'win32':
            return _try_win32com(trusted_copy, dest)
        else:
            # macOS / Linux: usa docx2pdf (AppleScript) con timeout
            import concurrent.futures
            try:
                from docx2pdf import convert as _convert
                executor = concurrent.futures.ThreadPoolExecutor(max_workers=1)
                future = executor.submit(_convert, trusted_copy, dest)
                try:
                    future.result(timeout=20)
                except concurrent.futures.TimeoutError:
                    executor.shutdown(wait=False)  # non bloccare sul thread zombie
                    return None, 'timeout docx2pdf (20s) — file bloccato in Word'
                finally:
                    executor.shutdown(wait=False)
                if os.path.isfile(dest) and os.path.getsize(dest) > 0:
                    return dest, None
                return None, 'output PDF vuoto'
            except Exception as e:
                return None, str(e)
    except Exception as e:
        return None, str(e)
    finally:
        shutil.rmtree(tmp_trusted, ignore_errors=True)


# ── DOCX / DOC / RTF ─────────────────────────────────────────────────────────

def _convert_docx_to_pdf(file_path: str, output_dir: str, lo_path: str) -> str:
    ext = os.path.splitext(file_path)[1].lower()
    step_errors = []

    # 1) Office COM (Word/Mac Word via docx2pdf)
    result, err = _try_docx2pdf(file_path, output_dir)
    if result:
        return result
    step_errors.append(f'Office COM: {err or "?"}')

    # 1b) Fallback AppleScript diretto — solo per .doc binario (Word 97-2003)
    # Per .docx/.rtf docx2pdf ha già tentato Word: inutile riprovare con AppleScript
    if sys.platform == 'darwin' and ext == '.doc':
        result, err = _try_applescript_word(file_path, output_dir)
        if result:
            return result
        step_errors.append(f'AppleScript Word: {err or "?"}')

    # 2) mammoth -> testo -> fpdf2 (non richiede GTK/WeasyPrint)
    if ext in ('.docx', '.doc', '.rtf'):
        try:
            import mammoth
            import zipfile
            from fpdf import FPDF

            # Valida che il file sia un ZIP valido prima di passarlo a mammoth
            # (un .docx corrotto può bloccare indefinitamente il parser ZIP)
            if ext in ('.docx', '.doc'):
                try:
                    with zipfile.ZipFile(file_path, 'r') as _z:
                        _z.testzip()
                except Exception as _ze:
                    raise RuntimeError(f'File non è un docx valido (ZIP corrotto): {_ze}')

            with open(file_path, 'rb') as f:
                raw = mammoth.extract_raw_text(f)

            text = raw.value
            if text.strip():
                base = os.path.splitext(os.path.basename(file_path))[0]
                dest = os.path.join(output_dir, f'{base}_{uuid.uuid4().hex[:8]}.pdf')

                pdf = FPDF()
                pdf.set_auto_page_break(auto=True, margin=15)
                pdf.add_page()
                pdf.set_font('Helvetica', size=11)

                for line in text.split('\n'):
                    try:
                        pdf.multi_cell(0, 6, line if line.strip() else ' ')
                    except Exception:
                        safe = line.encode('latin-1', errors='replace').decode('latin-1')
                        pdf.multi_cell(0, 6, safe if safe.strip() else ' ')

                pdf.output(dest)
                if os.path.isfile(dest) and os.path.getsize(dest) > 0:
                    return dest
                step_errors.append('mammoth+fpdf2: output vuoto')
            else:
                step_errors.append('mammoth: documento senza testo estraibile')
        except Exception as e:
            step_errors.append(f'mammoth+fpdf2: {e}')

    # 3) LibreOffice
    if lo_path:
        return _convert_office_to_pdf(file_path, output_dir, lo_path)

    raise RuntimeError(
        f'Impossibile convertire {os.path.basename(file_path)}: '
        + ' | '.join(step_errors)
    )


# ── XLSX / XLS / CSV / ODS ───────────────────────────────────────────────────

def _convert_xlsx_to_pdf(file_path: str, output_dir: str, lo_path: str) -> str:
    ext = os.path.splitext(file_path)[1].lower()
    step_errors = []

    # 1) Office COM (Excel, solo Windows/macOS con Office)
    result, err = _try_docx2pdf(file_path, output_dir)
    if result:
        return result
    step_errors.append(f'Office COM: {err or "?"}')

    # 2) openpyxl + fpdf2 (tabelle formattate, Python puro)
    if ext in ('.xlsx', '.xls', '.csv', '.ods'):
        try:
            from fpdf import FPDF

            # -- Lettura dati --
            if ext == '.csv':
                import csv
                with open(file_path, newline='', encoding='utf-8-sig', errors='replace') as f:
                    reader = csv.reader(f)
                    raw_rows = [r for r in reader if any(c.strip() for c in r)]
                sheets_data = {
                    os.path.splitext(os.path.basename(file_path))[0]: raw_rows
                }
            else:
                import openpyxl
                wb = openpyxl.load_workbook(file_path, data_only=True, read_only=True)
                sheets_data = {}
                for name in wb.sheetnames:
                    ws = wb[name]
                    sheet_rows = []
                    for row in ws.iter_rows(values_only=True):
                        cells = [str(c) if c is not None else '' for c in row]
                        if any(c.strip() for c in cells):
                            sheet_rows.append(cells)
                    if sheet_rows:
                        sheets_data[name] = sheet_rows
                wb.close()

            if not sheets_data:
                raise ValueError('Nessun dato trovato nel foglio')

            base = os.path.splitext(os.path.basename(file_path))[0]
            dest = os.path.join(output_dir, f'{base}_{uuid.uuid4().hex[:8]}.pdf')

            pdf = FPDF(orientation='L', unit='mm', format='A4')
            pdf.set_margins(10, 10, 10)
            pdf.set_auto_page_break(auto=True, margin=10)

            for sheet_name, sheet_rows in sheets_data.items():
                pdf.add_page()

                # Titolo foglio
                pdf.set_font('Helvetica', 'B', 12)
                pdf.set_fill_color(27, 107, 69)
                pdf.set_text_color(255, 255, 255)
                pdf.cell(0, 8, sheet_name, fill=True, ln=True)
                pdf.set_text_color(0, 0, 0)
                pdf.ln(2)

                col_count = max(len(r) for r in sheet_rows)
                if col_count == 0:
                    continue

                page_w = pdf.w - 20
                col_w = max(10, min(50, page_w / col_count))
                row_h = 5

                for i, row in enumerate(sheet_rows):
                    row = list(row) + [''] * (col_count - len(row))

                    if i == 0:
                        pdf.set_font('Helvetica', 'B', 7)
                        pdf.set_fill_color(220, 237, 228)
                    elif i % 2 == 0:
                        pdf.set_font('Helvetica', '', 7)
                        pdf.set_fill_color(255, 255, 255)
                    else:
                        pdf.set_font('Helvetica', '', 7)
                        pdf.set_fill_color(245, 247, 246)

                    for cell in row:
                        text = str(cell)
                        if len(text) > 35:
                            text = text[:33] + '\u2026'
                        pdf.cell(col_w, row_h, text, border=1, fill=True)
                    pdf.ln(row_h)

            pdf.output(dest)
            if os.path.isfile(dest) and os.path.getsize(dest) > 0:
                return dest
        except Exception:
            pass

    # 3) LibreOffice
    if lo_path:
        return _convert_office_to_pdf(file_path, output_dir, lo_path)

    raise RuntimeError(
        f'Impossibile convertire {os.path.basename(file_path)}: '
        + ' | '.join(step_errors)
    )


# ── PPTX / PPT / ODP ─────────────────────────────────────────────────────────

def _convert_pptx_to_pdf(file_path: str, output_dir: str, lo_path: str) -> str:
    ext = os.path.splitext(file_path)[1].lower()
    step_errors = []

    # 1) Office COM
    result, err = _try_docx2pdf(file_path, output_dir)
    if result:
        return result
    step_errors.append(f'Office COM: {err or "?"}')

    # 2) python-pptx -> testo per slide -> fpdf2
    if ext in ('.pptx', '.ppt'):
        try:
            from pptx import Presentation
            from fpdf import FPDF

            prs = Presentation(file_path)
            base = os.path.splitext(os.path.basename(file_path))[0]
            dest = os.path.join(output_dir, f'{base}_{uuid.uuid4().hex[:8]}.pdf')

            pdf = FPDF(orientation='L', unit='mm', format='A4')
            pdf.set_auto_page_break(auto=True, margin=15)

            for slide_num, slide in enumerate(prs.slides, 1):
                pdf.add_page()
                pdf.set_font('Helvetica', '', 8)
                pdf.set_text_color(150, 150, 150)
                pdf.cell(0, 5, f'Slide {slide_num}', ln=True, align='R')
                pdf.set_text_color(0, 0, 0)

                for shape in slide.shapes:
                    if not shape.has_text_frame:
                        continue
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if not text:
                            continue
                        ph = getattr(shape, 'placeholder_format', None)
                        is_title = ph is not None and getattr(ph, 'idx', -1) == 0
                        if is_title:
                            pdf.set_font('Helvetica', 'B', 16)
                            pdf.set_text_color(27, 107, 69)
                        else:
                            pdf.set_font('Helvetica', '', 11)
                            pdf.set_text_color(0, 0, 0)
                        pdf.multi_cell(0, 7, text)
                        pdf.ln(1)

            pdf.output(dest)
            if os.path.isfile(dest) and os.path.getsize(dest) > 0:
                return dest
        except Exception:
            pass

    # 3) LibreOffice
    if lo_path:
        return _convert_office_to_pdf(file_path, output_dir, lo_path)

    raise RuntimeError(
        f'Impossibile convertire {os.path.basename(file_path)}: '
        + ' | '.join(step_errors)
    )


# ── HTML ─────────────────────────────────────────────────────────────────────

def _convert_html_to_pdf(file_path: str, output_dir: str, lo_path: str) -> str:
    step_errors = []

    # 1) WeasyPrint (richiede GTK su Windows)
    try:
        import weasyprint
        base = os.path.splitext(os.path.basename(file_path))[0]
        dest = os.path.join(output_dir, f'{base}_{uuid.uuid4().hex[:8]}.pdf')
        weasyprint.HTML(filename=file_path).write_pdf(dest)
        if os.path.isfile(dest) and os.path.getsize(dest) > 0:
            return dest
        step_errors.append('weasyprint: output vuoto')
    except Exception as e:
        step_errors.append(f'weasyprint: {e}')

    # 2) Estrai testo via html.parser → fpdf2 (funziona senza GTK)
    try:
        from html.parser import HTMLParser
        from fpdf import FPDF

        class _TextExtractor(HTMLParser):
            def __init__(self):
                super().__init__()
                self.lines = []
                self._current = []
                self._skip = False
            def handle_starttag(self, tag, attrs):
                if tag in ('script', 'style'):
                    self._skip = True
                if tag in ('p', 'br', 'h1', 'h2', 'h3', 'li', 'div', 'tr'):
                    if self._current:
                        self.lines.append(''.join(self._current).strip())
                        self._current = []
            def handle_endtag(self, tag):
                if tag in ('script', 'style'):
                    self._skip = False
            def handle_data(self, data):
                if not self._skip:
                    self._current.append(data)
            def get_text(self):
                if self._current:
                    self.lines.append(''.join(self._current).strip())
                return '\n'.join(l for l in self.lines if l)

        with open(file_path, 'r', encoding='utf-8', errors='replace') as f:
            html_src = f.read()

        parser = _TextExtractor()
        parser.feed(html_src)
        text = parser.get_text()

        if text.strip():
            base = os.path.splitext(os.path.basename(file_path))[0]
            dest = os.path.join(output_dir, f'{base}_{uuid.uuid4().hex[:8]}.pdf')
            pdf = FPDF()
            pdf.set_auto_page_break(auto=True, margin=15)
            pdf.add_page()
            pdf.set_font('Helvetica', size=10)
            for line in text.split('\n'):
                try:
                    pdf.multi_cell(0, 6, line if line.strip() else ' ')
                except Exception:
                    safe = line.encode('latin-1', errors='replace').decode('latin-1')
                    pdf.multi_cell(0, 6, safe if safe.strip() else ' ')
            pdf.output(dest)
            if os.path.isfile(dest) and os.path.getsize(dest) > 0:
                return dest
        step_errors.append('html→fpdf2: testo vuoto o output mancante')
    except Exception as e:
        step_errors.append(f'html→fpdf2: {e}')

    # 3) LibreOffice
    if lo_path:
        return _convert_office_to_pdf(file_path, output_dir, lo_path)

    raise RuntimeError(
        f'Impossibile convertire HTML {os.path.basename(file_path)}: '
        + ' | '.join(step_errors)
    )


# ── TXT / XML ─────────────────────────────────────────────────────────────────

def _convert_txt_to_pdf(file_path: str, output_dir: str, lo_path: str) -> str:
    try:
        from fpdf import FPDF
        pdf = FPDF()
        pdf.set_auto_page_break(auto=True, margin=15)
        pdf.add_page()
        pdf.set_font('Helvetica', size=9)

        with open(file_path, 'r', encoding='utf-8', errors='replace') as f:
            content = f.read()

        for line in content.split('\n'):
            pdf.multi_cell(0, 5, line if line else ' ')

        base = os.path.splitext(os.path.basename(file_path))[0]
        dest = os.path.join(output_dir, f'{base}_{uuid.uuid4().hex[:8]}.pdf')
        pdf.output(dest)
        return dest
    except Exception:
        pass

    if lo_path:
        return _convert_office_to_pdf(file_path, output_dir, lo_path)

    raise RuntimeError(f'Impossibile convertire testo: {os.path.basename(file_path)}')


# ── Immagini ──────────────────────────────────────────────────────────────────

def _convert_image_to_pdf(file_path: str, output_dir: str) -> str:
    base = os.path.splitext(os.path.basename(file_path))[0]
    dest = os.path.join(output_dir, f'{base}_{uuid.uuid4().hex[:8]}.pdf')

    try:
        import img2pdf
        from PIL import Image

        with Image.open(file_path) as img:
            if img.mode in ('RGBA', 'LA', 'P'):
                conv = os.path.join(output_dir, f'_img_{uuid.uuid4().hex[:6]}.jpg')
                img.convert('RGB').save(conv, 'JPEG', quality=95)
                with open(dest, 'wb') as f:
                    f.write(img2pdf.convert(conv))
                os.unlink(conv)
            else:
                with open(dest, 'wb') as f:
                    f.write(img2pdf.convert(file_path))
        return dest
    except Exception:
        pass

    try:
        from PIL import Image
        with Image.open(file_path) as img:
            img.convert('RGB').save(dest, 'PDF', resolution=150)
        return dest
    except Exception as e:
        raise RuntimeError(f'Impossibile convertire immagine: {e}')


# ── Entry point pubblico ──────────────────────────────────────────────────────

def convert_to_pdf(file_path: str, output_dir: str, lo_path: str = None) -> str:
    """
    Converte qualsiasi file supportato in PDF.
    lo_path e' opzionale: se None, usa solo le librerie Python.
    """
    ext = os.path.splitext(file_path)[1].lower()

    # P7M: estrai il contenuto e converti ricorsivamente
    if ext == '.p7m':
        from core.p7m_handler import extract_p7m
        extracted = extract_p7m(file_path, output_dir)
        if extracted is None:
            raise RuntimeError('Impossibile estrarre il contenuto dal file P7M')
        try:
            return convert_to_pdf(extracted, output_dir, lo_path)
        finally:
            if os.path.exists(extracted):
                os.unlink(extracted)

    # PDF: copia diretta
    if ext == '.pdf':
        base = os.path.splitext(os.path.basename(file_path))[0]
        dest = os.path.join(output_dir, f'{base}_{uuid.uuid4().hex[:8]}.pdf')
        shutil.copy2(file_path, dest)
        return dest

    # Immagini
    if ext in IMAGE_EXTENSIONS:
        return _convert_image_to_pdf(file_path, output_dir)

    # Testo / XML semplice
    if ext in TEXT_EXTENSIONS or ext in XML_EXTENSIONS:
        return _convert_txt_to_pdf(file_path, output_dir, lo_path)

    # HTML
    if ext in HTML_EXTENSIONS:
        return _convert_html_to_pdf(file_path, output_dir, lo_path)

    # DOCX / DOC / RTF
    if ext in DOCX_EXTENSIONS:
        return _convert_docx_to_pdf(file_path, output_dir, lo_path)

    # XLSX / XLS / CSV / ODS
    if ext in XLSX_EXTENSIONS:
        return _convert_xlsx_to_pdf(file_path, output_dir, lo_path)

    # PPTX / PPT / ODP
    if ext in PPTX_EXTENSIONS:
        return _convert_pptx_to_pdf(file_path, output_dir, lo_path)

    # ODT / ODG - formati nativi LibreOffice
    if ext in ODT_EXTENSIONS:
        if lo_path:
            return _convert_office_to_pdf(file_path, output_dir, lo_path)
        # Tentativo mammoth per ODT
        try:
            import mammoth, weasyprint
            with open(file_path, 'rb') as f:
                doc = mammoth.convert_to_html(f)
            base = os.path.splitext(os.path.basename(file_path))[0]
            dest = os.path.join(output_dir, f'{base}_{uuid.uuid4().hex[:8]}.pdf')
            weasyprint.HTML(string=f'<html><body>{doc.value}</body></html>').write_pdf(dest)
            return dest
        except Exception:
            pass
        raise RuntimeError(
            f'{ext} richiede LibreOffice (scaricalo da libreoffice.org)'
        )

    # Fallback generico
    if lo_path:
        return _convert_office_to_pdf(file_path, output_dir, lo_path)

    raise RuntimeError(f'Formato non supportato: {ext}')
